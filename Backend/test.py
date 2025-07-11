"""
Local ID Card Generator for Sample Creation
Generates 20 random student ID cards and saves them locally to your device
No Google Drive upload - just for showing samples of how ID cards look
"""

import os
import io
import random
import requests
from datetime import datetime, timedelta
from faker import Faker
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from PIL import Image

# Initialize Faker for Indian data
fake = Faker('en_IN')

class LocalIDCardGenerator:
    """Generate ID cards locally for samples"""
    
    def __init__(self, template_path='./Jain.pptx', output_folder='./sample_id_cards'):
        """
        Initialize the local ID card generator
        
        Args:
            template_path (str): Path to your PowerPoint template
            output_folder (str): Folder to save generated ID cards
        """
        self.template_path = template_path
        self.output_folder = output_folder
        
        # Create output folder
        os.makedirs(self.output_folder, exist_ok=True)
        
        # Sample data for generation
        self.departments = [
            'Computer Science Engineering',
            'Electronics and Communication Engineering', 
            'Mechanical Engineering',
            'Civil Engineering',
            'Information Technology',
            'Electrical Engineering',
            'Biotechnology',
            'Chemical Engineering',
            'Aerospace Engineering',
            'Artificial Intelligence and Machine Learning'
        ]
        
        self.states = [
            'Karnataka', 'Tamil Nadu', 'Kerala', 'Andhra Pradesh', 
            'Maharashtra', 'Gujarat', 'Rajasthan', 'Uttar Pradesh',
            'West Bengal', 'Bihar', 'Odisha', 'Madhya Pradesh'
        ]
        
        self.blood_groups = ['A+', 'A-', 'B+', 'B-', 'AB+', 'AB-', 'O+', 'O-']
        self.categories = ['General', 'OBC', 'SC', 'ST', 'EWS']
        self.religions = ['Hindu', 'Muslim', 'Christian', 'Sikh', 'Buddhist', 'Jain']
        
    def generate_student_data(self, index):
        """Generate random student data"""
        gender = random.choice(['Male', 'Female'])
        first_name = fake.first_name_male() if gender == 'Male' else fake.first_name_female()
        last_name = fake.last_name()
        full_name = f"{first_name} {last_name}"
        
        birth_date = fake.date_of_birth(minimum_age=18, maximum_age=23)
        department = random.choice(self.departments)
        state = random.choice(self.states)
        
        # Generate marks
        tenth_total = 500
        tenth_scored = random.randint(400, 485)
        tenth_percentage = round((tenth_scored / tenth_total) * 100, 2)
        
        twelfth_total = 500
        twelfth_scored = random.randint(420, 490)
        twelfth_percentage = round((twelfth_scored / twelfth_total) * 100, 2)
        
        # Subject marks
        physics_scored = random.randint(80, 98)
        chemistry_scored = random.randint(80, 98)
        mathematics_scored = random.randint(80, 98)
        pcm_total = physics_scored + chemistry_scored + mathematics_scored
        pcm_percentage = round((pcm_total / 300) * 100, 2)
        
        return {
            'student_id': f"STU{random.randint(10000000, 99999999)}",
            'studentFullName': full_name,
            'name': full_name,
            'gender': gender,
            'dateOfBirth': birth_date.strftime('%Y-%m-%d'),
            'dob': birth_date.strftime('%Y-%m-%d'),
            'bloodGroup': random.choice(self.blood_groups),
            'nationality': 'Indian',
            'religion': random.choice(self.religions),
            'caste': fake.last_name(),
            'motherTongue': random.choice(['Hindi', 'Telugu', 'Tamil', 'Kannada', 'Malayalam']),
            'category': random.choice(self.categories),
            'birthPlace': fake.city(),
            
            # Contact Information
            'studentContactNo': f"+91{random.randint(7000000000, 9999999999)}",
            'phone': f"+91{random.randint(7000000000, 9999999999)}",
            'studentEmail': f"{first_name.lower()}.{last_name.lower()}{index}@college.edu",
            'email': f"{first_name.lower()}.{last_name.lower()}{index}@college.edu",
            'parentContactNo': f"+91{random.randint(7000000000, 9999999999)}",
            'parentEmail': f"parent{index}@email.com",
            'parent_email': f"parent{index}@email.com",
            
            # Address
            'correspondenceAddress': fake.address(),
            'correspondenceCity': fake.city(),
            'correspondenceState': state,
            'correspondenceCountry': 'India',
            'correspondencePostalCode': fake.postcode(),
            
            # Academic Information
            'tenthBoardUniversity': random.choice(['CBSE', 'ICSE', 'State Board']),
            'tenthSchoolName': f"{fake.city()} High School",
            'tenthSchoolState': state,
            'tenthPassedOutYear': str(birth_date.year + 16),
            'tenthTotalMarks': str(tenth_total),
            'tenthScoredMarks': str(tenth_scored),
            'tenthPercentage': str(tenth_percentage),
            
            'boardUniversity': random.choice(['CBSE', 'ICSE', 'State Board']),
            'collegeInstitutionName': f"{fake.city()} College",
            'collegeStateName': state,
            'twelfthPassedOutYear': str(birth_date.year + 18),
            'twelfthTotalMarks': str(twelfth_total),
            'twelfthScoredMarks': str(twelfth_scored),
            'twelfthPercentage': str(twelfth_percentage),
            
            # Subject marks
            'pcmPercentage': str(pcm_percentage),
            'physicsTotal': '100',
            'physicsScored': str(physics_scored),
            'chemistryTotal': '100',
            'chemistryScored': str(chemistry_scored),
            'mathematicsTotal': '100',
            'mathematicsScored': str(mathematics_scored),
            
            # Parent Information
            'fatherName': fake.name_male(),
            'fatherOccupation': random.choice(['Engineer', 'Teacher', 'Businessman', 'Doctor', 'Farmer']),
            'fatherIncome': str(random.randint(300000, 1500000)),
            'fatherMobile': f"+91{random.randint(7000000000, 9999999999)}",
            'motherName': fake.name_female(),
            'motherOccupation': random.choice(['Teacher', 'Housewife', 'Nurse', 'Engineer', 'Doctor']),
            'motherMobile': f"+91{random.randint(7000000000, 9999999999)}",
            
            # University Details
            'department': department,
            'programName': 'B.Tech',
            'course': 'B.Tech',
            'admissionType': random.choice(['Regular', 'Management', 'Sports']),
            'juApplication': f"JU{datetime.now().year}{random.randint(100000, 999999)}",
            
            # Dates
            'issue_date': datetime.now().strftime('%Y-%m-%d'),
            'expiry_date': (datetime.now() + timedelta(days=4*365)).strftime('%Y-%m-%d'),
            'registration_date': fake.date_between(start_date='-6m', end_date='today').strftime('%Y-%m-%d'),
            
            # Photo URL (random photo)
            'photo_url': f"https://randomuser.me/api/portraits/{'men' if gender == 'Male' else 'women'}/{random.randint(1, 99)}.jpg"
        }
    
    def download_photo(self, photo_url):
        """Download photo and return as BytesIO stream"""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            response = requests.get(photo_url, timeout=10, headers=headers)
            if response.status_code == 200:
                image = Image.open(io.BytesIO(response.content))
                if image.mode != 'RGB':
                    image = image.convert('RGB')
                
                # Resize to appropriate size
                image.thumbnail((300, 400), Image.Resampling.LANCZOS)
                
                # Convert to BytesIO
                img_stream = io.BytesIO()
                image.save(img_stream, format='JPEG', quality=85)
                img_stream.seek(0)
                return img_stream
        except Exception as e:
            print(f"Failed to download photo: {e}")
        return None
    
    def find_photo_placeholder(self, slide):
        """Find photo placeholder in slide"""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content = shape.text.upper()
                if 'PHOTO' in text_content or '{{PHOTO}}' in text_content:
                    return shape
            if hasattr(shape, 'name') and 'photo' in shape.name.lower():
                return shape
        return None
    
    def replace_text_placeholders(self, shape_or_cell, student_data):
        """Replace text placeholders with student data"""
        if not hasattr(shape_or_cell, 'text'):
            return
        
        original_text = shape_or_cell.text
        if not original_text:
            return
        
        # Address formatting
        address_lines = []
        if student_data.get('correspondenceAddress'):
            address_lines.append(student_data['correspondenceAddress'])
        if student_data.get('correspondenceCity'):
            address_lines.append(student_data['correspondenceCity'])
        if student_data.get('correspondenceState') and student_data.get('correspondencePostalCode'):
            address_lines.append(f"{student_data['correspondenceState']} - {student_data['correspondencePostalCode']}")
        address = '\n'.join(address_lines)
        
        # Placeholder mappings
        replacements = {
            '{{NAME}}': str(student_data.get('name', student_data.get('studentFullName', ''))),
            '{{ID_NUMBER}}': str(student_data.get('juApplication', '')),
            '{{DEPARTMENT}}': str(student_data.get('department', '')),
            '{{COURSE}}': str(student_data.get('course', student_data.get('programName', ''))),
            '{{PROGRAM}}': str(student_data.get('programName', '')),
            '{{ADDRESS}}': address,
            '{{MOBILE}}': str(student_data.get('phone', student_data.get('studentContactNo', ''))),
            '{{PHONE}}': str(student_data.get('studentContactNo', '')),
            '{{EMAIL}}': str(student_data.get('email', student_data.get('studentEmail', ''))),
            '{{EMERGENCY_CONTACT}}': str(student_data.get('parentContactNo', '')),
            '{{PARENT_EMAIL}}': str(student_data.get('parentEmail', student_data.get('parent_email', ''))),
            '{{BLOOD_GROUP}}': str(student_data.get('bloodGroup', '')),
            '{{DATE_OF_BIRTH}}': str(student_data.get('dateOfBirth', student_data.get('dob', ''))),
            '{{DOB}}': str(student_data.get('dateOfBirth', student_data.get('dob', ''))),
            '{{FATHER_NAME}}': str(student_data.get('fatherName', '')),
            '{{FATHER_OCCUPATION}}': str(student_data.get('fatherOccupation', '')),
            '{{FATHER_MOBILE}}': str(student_data.get('fatherMobile', '')),
            '{{MOTHER_NAME}}': str(student_data.get('motherName', '')),
            '{{MOTHER_OCCUPATION}}': str(student_data.get('motherOccupation', '')),
            '{{MOTHER_MOBILE}}': str(student_data.get('motherMobile', '')),
            '{{TENTH_BOARD}}': str(student_data.get('tenthBoardUniversity', '')),
            '{{TENTH_SCHOOL}}': str(student_data.get('tenthSchoolName', '')),
            '{{TENTH_YEAR}}': str(student_data.get('tenthPassedOutYear', '')),
            '{{TENTH_PERCENTAGE}}': str(student_data.get('tenthPercentage', '')),
            '{{TWELFTH_BOARD}}': str(student_data.get('boardUniversity', '')),
            '{{TWELFTH_COLLEGE}}': str(student_data.get('collegeInstitutionName', '')),
            '{{TWELFTH_YEAR}}': str(student_data.get('twelfthPassedOutYear', '')),
            '{{TWELFTH_PERCENTAGE}}': str(student_data.get('twelfthPercentage', '')),
            '{{PCM_PERCENTAGE}}': str(student_data.get('pcmPercentage', '')),
            '{{ISSUE_DATE}}': str(student_data.get('issue_date', '')),
            '{{EXPIRY_DATE}}': str(student_data.get('expiry_date', '')),
            '{{REGISTRATION_DATE}}': str(student_data.get('registration_date', '')),
            '{{PHOTO}}': '',  # Handled separately
            '{{BARCODE}}': str(student_data.get('student_id', '')),
            '{{QR_CODE}}': str(student_data.get('student_id', '')),
        }
        
        # Replace placeholders
        new_text = original_text
        for placeholder, value in replacements.items():
            if placeholder in new_text:
                new_text = new_text.replace(placeholder, value)
        
        if new_text != original_text:
            shape_or_cell.text = new_text
            
            # Apply formatting
            if hasattr(shape_or_cell, 'text_frame') and shape_or_cell.text_frame:
                for paragraph in shape_or_cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Red Hat Display Bold'
                        run.font.size = Pt(8)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)
    
    def generate_single_id_card(self, student_data):
        """Generate a single ID card"""
        try:
            if not os.path.exists(self.template_path):
                return {'success': False, 'error': f'Template not found: {self.template_path}'}
            
            # Load template
            ppt = Presentation(self.template_path)
            
            # Process all slides
            for slide_num, slide in enumerate(ppt.slides):
                # Replace text in all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        self.replace_text_placeholders(shape, student_data)
                    
                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                self.replace_text_placeholders(cell, student_data)
                
                # Insert photo on first slide
                if slide_num == 0:
                    photo_shape = self.find_photo_placeholder(slide)
                    if photo_shape and student_data.get('photo_url'):
                        try:
                            photo_stream = self.download_photo(student_data['photo_url'])
                            if photo_stream:
                                # Get placeholder position and size
                                left = photo_shape.left
                                top = photo_shape.top
                                width = photo_shape.width
                                height = photo_shape.height
                                
                                # Remove placeholder
                                slide.shapes._spTree.remove(photo_shape._element)
                                
                                # Add photo
                                slide.shapes.add_picture(photo_stream, left, top, width, height)
                        except Exception as e:
                            print(f"Failed to insert photo: {e}")
            
            # Save file
            safe_name = ''.join(c for c in student_data['studentFullName'] if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_name = safe_name.replace(' ', '_')[:50]
            filename = f"ID_Card_{safe_name}_{student_data['student_id']}.pptx"
            filepath = os.path.join(self.output_folder, filename)
            
            ppt.save(filepath)
            
            return {
                'success': True,
                'filename': filename,
                'filepath': filepath,
                'student_name': student_data['studentFullName'],
                'student_id': student_data['student_id']
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'student_name': student_data.get('studentFullName', 'Unknown')
            }
    
    def generate_batch_id_cards(self, count=20):
        """Generate multiple ID cards"""
        print(f"🎓 Generating {count} sample ID cards...")
        print(f"📁 Output folder: {os.path.abspath(self.output_folder)}")
        print("=" * 60)
        
        results = []
        successful = 0
        
        for i in range(count):
            print(f"\n📋 Generating ID Card {i+1}/{count}")
            
            # Generate student data
            student_data = self.generate_student_data(i+1)
            print(f"Student: {student_data['studentFullName']} ({student_data['department']})")
            
            # Generate ID card
            result = self.generate_single_id_card(student_data)
            
            if result['success']:
                print(f"✅ ID card saved: {result['filename']}")
                successful += 1
            else:
                print(f"❌ Failed: {result['error']}")
            
            results.append(result)
        
        print(f"\n" + "=" * 60)
        print(f"📊 GENERATION COMPLETE")
        print(f"✅ Successfully generated: {successful}/{count} ID cards")
        print(f"❌ Failed: {count - successful}/{count} ID cards")
        print(f"📁 All files saved in: {os.path.abspath(self.output_folder)}")
        
        # List generated files
        if successful > 0:
            print(f"\n📄 Generated Files:")
            for result in results:
                if result['success']:
                    print(f"• {result['filename']}")
        
        return results


def main():
    """Main function"""
    print("🎓 Local ID Card Sample Generator")
    print("Generates 20 sample ID cards and saves them to your device")
    print("=" * 60)
    
    # Check for template file
    template_path = './Jain.pptx'
    if not os.path.exists(template_path):
        template_path = input("Enter path to your PowerPoint template (Jain.pptx): ").strip()
        if not os.path.exists(template_path):
            print(f"❌ Template file not found: {template_path}")
            print("Please make sure your Jain.pptx template is in the same folder as this script.")
            return
    
    # Set output folder
    output_folder = './sample_id_cards'
    custom_folder = input(f"Enter output folder (or press Enter for '{output_folder}'): ").strip()
    if custom_folder:
        output_folder = custom_folder
    
    # Initialize generator
    generator = LocalIDCardGenerator(template_path, output_folder)
    
    # Generate ID cards
    print(f"\n🚀 Starting generation...")
    results = generator.generate_batch_id_cards(20)
    
    print(f"\n🎉 Generation complete!")
    print(f"Check the '{output_folder}' folder for your sample ID cards.")
    
    # Open folder (Windows/Mac)
    try:
        import subprocess
        import platform
        
        if platform.system() == "Windows":
            subprocess.run(['explorer', os.path.abspath(output_folder)])
        elif platform.system() == "Darwin":  # macOS
            subprocess.run(['open', os.path.abspath(output_folder)])
        else:  # Linux
            subprocess.run(['xdg-open', os.path.abspath(output_folder)])
    except:
        pass


if __name__ == "__main__":
    main()