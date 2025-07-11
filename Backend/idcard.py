"""
Updated Student ID Card Generator Module with Google Drive Upload Support
Generates ID cards and uploads them to the student's Google Drive folder.
"""

import os
import io
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from datetime import datetime
from PIL import Image

# Google Drive imports
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google.oauth2.service_account import Credentials

# Windows-specific imports for PDF conversion
try:
    import comtypes.client
    WINDOWS_AVAILABLE = True
except ImportError:
    WINDOWS_AVAILABLE = False
    print("Warning: PDF conversion requires Windows + PowerPoint")


class StudentIDCardGenerator:
    """
    Student ID Card Generator class for creating formatted ID cards with Google Drive upload.
    """
    
    def __init__(self, template_path='./Jain.pptx', 
                 output_folder='./generated_id_cards',
                 service_account_file='./credentials.json'):
        """
        Initialize the ID card generator.
        
        Args:
            template_path (str): Path to PowerPoint template
            output_folder (str): Folder to save generated files
            service_account_file (str): Path to Google service account credentials
        """
        self.template_path = template_path
        self.output_folder = output_folder
        self.service_account_file = service_account_file
        
        # Create output folder if it doesn't exist
        os.makedirs(self.output_folder, exist_ok=True)
        
        # Initialize Google Drive service
        self.drive_service = self._get_drive_service()
    
    def _get_drive_service(self):
        """Create and return Google Drive API service"""
        try:
            scopes = ['https://www.googleapis.com/auth/drive']
            credentials = Credentials.from_service_account_file(
                self.service_account_file, scopes=scopes
            )
            service = build('drive', 'v3', credentials=credentials)
            return service
        except Exception as e:
            print(f"Error creating Drive service: {e}")
            return None
    
    def generate_id_card(self, student_data):
        """
        Generate ID card for a single student and upload to Google Drive.
        
        Args:
            student_data (dict): Student information dictionary
        
        Returns:
            dict: Result with success status and file paths (local and Google Drive)
        """
        try:
            # Validate template exists
            if not os.path.exists(self.template_path):
                return {
                    "success": False,
                    "error": f"ID card template not found: {self.template_path}"
                }
            
            student_name = student_data.get('studentFullName', '') or student_data.get('name', 'Unknown Student')
            print(f"Generating ID card for {student_name}...")
            
            # Load fresh template
            ppt = Presentation(self.template_path)
            
            # Replace placeholders in ALL slides (front and back pages)
            for slide_num, slide in enumerate(ppt.slides):
                print(f"  Processing slide {slide_num + 1}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        # Replace text placeholders with formatting
                        self._replace_text_with_formatting(shape, student_data)
                    
                    # Handle text in tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                self._replace_text_with_formatting(cell, student_data)
                
                # Handle photo insertion for first slide only (front of card)
                if slide_num == 0:
                    self._insert_student_photo(slide, student_data)
            
            # Generate safe filename
            safe_name = self._sanitize_filename(student_name)
            
            ppt_file = os.path.join(self.output_folder, f"ID_Card_{safe_name}.pptx")
            pdf_file = os.path.join(self.output_folder, f"ID_Card_{safe_name}.pdf")
            
            # Save PowerPoint
            ppt.save(ppt_file)
            print(f"  Generated: {ppt_file}")
            
            # Convert to PDF
            pdf_created = self._convert_to_pdf(ppt_file, pdf_file)
            if pdf_created:
                print(f"  PDF created: {pdf_file}")
            
            # Upload to Google Drive
            gdrive_result = self._upload_to_google_drive(student_data, ppt_file, pdf_file if pdf_created else None)
            
            # Determine primary file
            primary_file = pdf_file if pdf_created else ppt_file
            primary_format = "pdf" if pdf_created else "pptx"
            
            return {
                "success": True,
                "file_path": primary_file,
                "ppt_path": ppt_file,
                "pdf_path": pdf_file if pdf_created else None,
                "format": primary_format,
                "google_drive": gdrive_result,
                "message": f"ID card generated successfully for {student_name}"
            }
                
        except Exception as e:
            print(f"Error generating ID card: {e}")
            return {
                "success": False,
                "error": f"Error generating ID card: {str(e)}"
            }
    
    def _upload_to_google_drive(self, student_data, ppt_file, pdf_file=None):
        """
        Upload generated ID card files to student's Google Drive folder.
        
        Args:
            student_data (dict): Student information
            ppt_file (str): Path to PowerPoint file
            pdf_file (str): Path to PDF file (optional)
            
        Returns:
            dict: Upload results
        """
        if not self.drive_service:
            return {
                "success": False,
                "error": "Google Drive service not available"
            }
        
        try:
            # Get student's Google Drive folder ID
            folder_id = student_data.get('documentsFolder')
            if not folder_id:
                return {
                    "success": False,
                    "error": "Student's Google Drive folder not found"
                }
            
            print(f"  Uploading ID card to Google Drive folder: {folder_id}")
            
            uploaded_files = {}
            
            # Upload PowerPoint file
            if os.path.exists(ppt_file):
                ppt_result = self._upload_single_file_to_drive(
                    ppt_file, 
                    folder_id, 
                    "ID_Card.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                if ppt_result['success']:
                    uploaded_files['powerpoint'] = ppt_result
                    print(f"    PowerPoint uploaded: {ppt_result['file_id']}")
            
            # Upload PDF file if available
            if pdf_file and os.path.exists(pdf_file):
                pdf_result = self._upload_single_file_to_drive(
                    pdf_file, 
                    folder_id, 
                    "ID_Card.pdf",
                    "application/pdf"
                )
                if pdf_result['success']:
                    uploaded_files['pdf'] = pdf_result
                    print(f"    PDF uploaded: {pdf_result['file_id']}")
            
            return {
                "success": True,
                "folder_id": folder_id,
                "uploaded_files": uploaded_files,
                "total_uploaded": len(uploaded_files)
            }
            
        except Exception as e:
            print(f"  Error uploading to Google Drive: {e}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def _upload_single_file_to_drive(self, file_path, folder_id, drive_filename, mime_type):
        """
        Upload a single file to Google Drive.
        
        Args:
            file_path (str): Local file path
            folder_id (str): Google Drive folder ID
            drive_filename (str): Filename in Google Drive
            mime_type (str): MIME type of the file
            
        Returns:
            dict: Upload result
        """
        try:
            # Check if file already exists and delete it
            existing_files = self.drive_service.files().list(
                q=f"name='{drive_filename}' and parents in '{folder_id}'",
                fields="files(id, name)"
            ).execute()
            
            for existing_file in existing_files.get('files', []):
                self.drive_service.files().delete(fileId=existing_file['id']).execute()
                print(f"    Deleted existing file: {existing_file['name']}")
            
            # File metadata
            file_metadata = {
                'name': drive_filename,
                'parents': [folder_id],
                'description': f"Student ID Card generated on {datetime.now().isoformat()}"
            }
            
            # Create media upload object
            media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
            
            # Upload file
            uploaded_file = self.drive_service.files().create(
                body=file_metadata,
                media_body=media,
                fields='id,name,webViewLink,webContentLink,size'
            ).execute()
            
            # Make file publicly viewable
            try:
                permission = {
                    'role': 'reader',
                    'type': 'anyone'
                }
                self.drive_service.permissions().create(
                    fileId=uploaded_file.get('id'),
                    body=permission
                ).execute()
            except Exception as perm_error:
                print(f"    Permission setting warning: {perm_error}")
            
            return {
                'success': True,
                'file_id': uploaded_file.get('id'),
                'file_name': uploaded_file.get('name'),
                'web_view_link': uploaded_file.get('webViewLink'),
                'download_link': uploaded_file.get('webContentLink'),
                'file_size': uploaded_file.get('size')
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }
    
    def _get_google_drive_photo_url(self, student_data):
        """
        Extract Google Drive photo URL from student data.
        
        Args:
            student_data (dict): Student information dictionary
            
        Returns:
            str: Photo download URL or None
        """
        try:
            # Check if uploadedDocuments exists and has photographUpload
            uploaded_docs = student_data.get('uploadedDocuments', {})
            photo_upload = uploaded_docs.get('photographUpload', {})
            
            if photo_upload:
                # Get the download link
                download_link = photo_upload.get('download_link', '')
                file_id = photo_upload.get('file_id', '')
                
                print(f"  Found photo in Google Drive:")
                print(f"    File ID: {file_id}")
                print(f"    Download link: {download_link}")
                
                # Return the download link (best for direct download)
                if download_link:
                    return download_link
                
                # Fallback: construct download URL from file_id
                if file_id:
                    return f"https://drive.google.com/uc?id={file_id}&export=download"
            
            # Check legacy fields
            photo_url = student_data.get('photo_url', '')
            if photo_url:
                print(f"  Found photo URL: {photo_url}")
                return photo_url
            
            print("  No photo URL found in student data")
            return None
            
        except Exception as e:
            print(f"  Error extracting photo URL: {e}")
            return None
    
    def _insert_student_photo(self, slide, student_data):
        """
        Insert student photo into the slide.
        
        Args:
            slide: PowerPoint slide object
            student_data (dict): Dictionary with student information
        """
        try:
            # Get photo URL from Google Drive data
            photo_url = self._get_google_drive_photo_url(student_data)
            
            if not photo_url:
                print("  No photo URL available")
                return
            
            # Download and process the image
            image_stream = self._get_image_stream(photo_url)
            if not image_stream:
                print("  Failed to download image")
                return
            
            # Look for photo placeholder or shape with "PHOTO" text
            photo_shape = self._find_photo_placeholder(slide)
            
            if photo_shape:
                # Get position and size of the placeholder
                left = photo_shape.left
                top = photo_shape.top
                width = photo_shape.width
                height = photo_shape.height
                print(f"    Photo placeholder found at: {left}, {top}, {width}, {height}")
                
                # Remove the placeholder shape
                slide.shapes._spTree.remove(photo_shape._element)
                
                # Add the actual photo
                slide.shapes.add_picture(image_stream, left, top, width, height)
                print("  Photo inserted successfully")
            else:
                # If no placeholder found, add photo at default position
                left = Inches(0.3)
                top = Inches(2.35)
                width = Inches(0.7)
                height = Inches(0.9)
                
                slide.shapes.add_picture(image_stream, left, top, width, height)
                print("  Photo inserted at default position")
                
        except Exception as e:
            print(f"  Error inserting photo: {e}")
    
    def _find_photo_placeholder(self, slide):
        """
        Find the photo placeholder shape in the slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Shape object or None
        """
        for shape in slide.shapes:
            # Check if shape has text containing "PHOTO" or "{{PHOTO}}"
            if hasattr(shape, 'text'):
                text_content = shape.text.upper()
                if 'PHOTO' in text_content or '{{PHOTO}}' in text_content:
                    return shape
            
            # Check if shape name contains "photo"
            if hasattr(shape, 'name') and 'photo' in shape.name.lower():
                return shape
        
        return None
    
    def _get_image_stream(self, image_url):
        """
        Download image from Google Drive URL and convert to BytesIO stream.
        
        Args:
            image_url (str): Google Drive download URL
            
        Returns:
            BytesIO stream or None
        """
        try:
            print(f"  Downloading image from: {image_url}")
            
            # Set headers to mimic a browser request
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            
            # Download from URL with timeout
            response = requests.get(image_url, timeout=30, headers=headers, stream=True)
            response.raise_for_status()
            
            # Read the image data
            image_data = response.content
            print(f"  Downloaded {len(image_data)} bytes")
            
            # Process the image
            image = Image.open(io.BytesIO(image_data))
            
            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Resize image if too large (maintain aspect ratio)
            max_size = (400, 500)  # Maximum width, height
            image.thumbnail(max_size, Image.Resampling.LANCZOS)
            print(f"  Resized image to: {image.size}")
            
            # Save to BytesIO stream
            image_stream = io.BytesIO()
            image.save(image_stream, format='JPEG', quality=85)
            image_stream.seek(0)
            
            return image_stream
            
        except requests.RequestException as e:
            print(f"  Network error downloading image: {e}")
            return None
        except Exception as e:
            print(f"  Error processing image: {e}")
            return None
    
    def _replace_text_with_formatting(self, shape_or_cell, student_data):
        """
        Replace placeholders in text with student data and apply formatting.
        
        Args:
            shape_or_cell: PowerPoint shape or table cell
            student_data (dict): Dictionary with student information
        """
        if not hasattr(shape_or_cell, 'text'):
            return
        
        original_text = shape_or_cell.text
        if not original_text:
            return
        
        # Get placeholder mappings
        replacements = self._get_placeholder_mappings(student_data)
        
        # Check if any placeholder exists in the text
        has_placeholder = any(placeholder in original_text for placeholder in replacements.keys())
        
        if has_placeholder:
            # Replace all placeholders
            new_text = original_text
            for placeholder, value in replacements.items():
                if placeholder in new_text:
                    new_text = new_text.replace(placeholder, value)
            
            # Update the text
            shape_or_cell.text = new_text
            
            # Apply uniform formatting to all text
            if hasattr(shape_or_cell, 'text_frame') and shape_or_cell.text_frame:
                for paragraph in shape_or_cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Apply Red Hat Display Bold, size 8 to all runs
                        run.font.name = 'Red Hat Display Bold'
                        run.font.size = Pt(8)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black
    
    def _get_placeholder_mappings(self, student_data):
        """
        Get mapping of placeholders to actual student values.
        
        Args:
            student_data (dict): Dictionary with student information
        
        Returns:
            dict: Placeholder to value mappings
        """
        # Construct full address with state and pincode on same line
        address_line1 = student_data.get('correspondenceAddress', '').strip()
        address_line2 = student_data.get('correspondenceCity', '').strip()
        
        # Combine state and pincode on one line
        state = student_data.get('correspondenceState', '').strip()
        pincode = student_data.get('correspondencePostalCode', '').strip()
        state_pincode = f"{state} - {pincode}" if state and pincode else (state or pincode)
        
        # Build final address
        address_lines = []
        if address_line1:
            address_lines.append(address_line1)
        if address_line2:
            address_lines.append(address_line2)
        if state_pincode:
            address_lines.append(state_pincode)
        
        address = '\n'.join(address_lines)
        
        return {
            # Front page placeholders
            '{{NAME}}': str(student_data.get('name', '') or student_data.get('studentFullName', '')),
            '{{ID_NUMBER}}': str(student_data.get('juApplication', '')),
            '{{DEPARTMENT}}': str(student_data.get('department', '')),
            '{{COURSE}}': str(student_data.get('course', '') or student_data.get('programName', '')),
            '{{PROGRAM}}': str(student_data.get('programName', '')),
            
            # Contact Information
            '{{ADDRESS}}': address,
            '{{MOBILE}}': str(student_data.get('mobile', '') or student_data.get('studentContactNo', '')),
            '{{PHONE}}': str(student_data.get('studentContactNo', '')),
            '{{EMAIL}}': str(student_data.get('email', '') or student_data.get('studentEmail', '')),
            '{{EMERGENCY_CONTACT}}': str(student_data.get('parentContactNo', '')),
            '{{PARENT_EMAIL}}': str(student_data.get('parentEmail', '') or student_data.get('parent_email', '')),
            
            # Academic Information
            '{{BLOOD_GROUP}}': str(student_data.get('bloodGroup', '')),
            '{{DATE_OF_BIRTH}}': str(student_data.get('dateOfBirth', '') or student_data.get('dob', '')),
            '{{DOB}}': str(student_data.get('dateOfBirth', '') or student_data.get('dob', '')),
            
            # Family Information
            '{{FATHER_NAME}}': str(student_data.get('fatherName', '')),
            '{{FATHER_OCCUPATION}}': str(student_data.get('fatherOccupation', '')),
            '{{FATHER_INCOME}}': str(student_data.get('fatherIncome', '')),
            '{{FATHER_MOBILE}}': str(student_data.get('fatherMobile', '')),
            '{{MOTHER_NAME}}': str(student_data.get('motherName', '')),
            '{{MOTHER_OCCUPATION}}': str(student_data.get('motherOccupation', '')),
            '{{MOTHER_MOBILE}}': str(student_data.get('motherMobile', '')),
            '{{GUARDIAN_NAME}}': str(student_data.get('guardianName', '')),
            
            # 10th Class Information
            '{{TENTH_BOARD}}': str(student_data.get('tenthBoardUniversity', '')),
            '{{TENTH_SCHOOL}}': str(student_data.get('tenthSchoolName', '')),
            '{{TENTH_YEAR}}': str(student_data.get('tenthPassedOutYear', '')),
            '{{TENTH_PERCENTAGE}}': str(student_data.get('tenthPercentage', '')),
            
            # 12th Class Information
            '{{TWELFTH_BOARD}}': str(student_data.get('boardUniversity', '')),
            '{{TWELFTH_COLLEGE}}': str(student_data.get('collegeInstitutionName', '')),
            '{{TWELFTH_YEAR}}': str(student_data.get('twelfthPassedOutYear', '')),
            '{{TWELFTH_PERCENTAGE}}': str(student_data.get('twelfthPercentage', '')),
            '{{PCM_PERCENTAGE}}': str(student_data.get('pcmPercentage', '')),
            
            # Administrative placeholders
            '{{ISSUE_DATE}}': str(student_data.get('issue_date', datetime.now().strftime('%Y-%m-%d'))),
            '{{EXPIRY_DATE}}': str(student_data.get('expiry_date', '')),
            '{{REGISTRATION_DATE}}': str(student_data.get('registration_date', '')),
            
            # Photo placeholder (will be replaced with actual image)
            '{{PHOTO}}': '',  # This will be handled by image insertion
            '{{BARCODE}}': str(student_data.get('student_id', '')),
            '{{QR_CODE}}': str(student_data.get('student_id', '')),
        }
    
    def _sanitize_filename(self, name):
        """
        Create a safe filename from student name.
        
        Args:
            name (str): Student name
            
        Returns:
            str: Sanitized filename
        """
        if not name:
            return "Unknown_Student"
            
        # Replace problematic characters
        safe_chars = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789_-"
        sanitized = ''.join(c if c in safe_chars else '_' for c in name)
        
        # Remove multiple underscores and limit length
        while '__' in sanitized:
            sanitized = sanitized.replace('__', '_')
        
        return sanitized[:50]  # Limit filename length
    
    def _convert_to_pdf(self, ppt_path, pdf_path):
        """
        Convert PowerPoint to PDF.
        
        Args:
            ppt_path (str): Path to PowerPoint file
            pdf_path (str): Path for output PDF file
        
        Returns:
            bool: True if successful, False otherwise
        """
        if not WINDOWS_AVAILABLE:
            print("  PDF conversion requires Windows with PowerPoint installed.")
            return False
        
        try:
            # Initialize PowerPoint
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1
            
            # Open and convert
            presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
            presentation.ExportAsFixedFormat(
                os.path.abspath(pdf_path),
                ExportFormat=2,  # PDF format
                Intent=1,        # Print quality
                FrameSlides=0,
                HandoutOrder=1,
                OutputType=0,    # Full slides
                PrintHiddenSlides=0,
                PrintRange=None,
                RangeType=1,
                SlideShowName="",
                IncludeDocProps=1,
                KeepIRMSettings=1,
                DocStructureTags=1,
                BitmapMissingFonts=1,
                UseDocumentICCProfile=0
            )
            
            # Clean up
            presentation.Close()
            powerpoint.Quit()
            return True
            
        except Exception as e:
            print(f"  PDF conversion error: {e}")
            return False


class StudentDataFormatter:
    """
    Helper class to format student data for ID card generation.
    """
    
    @staticmethod
    def format_student_data(student_record, dept_record=None):
        """
        Format student data from database records.
        This handles the exact data structure you provided.
        
        Args:
            student_record (dict): Main student record from database
            dept_record (dict): Department-specific record (same as student_record in your case)
            
        Returns:
            dict: Formatted data for ID card template
        """
        # In your case, dept_record contains all the data we need
        data_source = dept_record if dept_record else student_record
        
        # Add calculated fields
        formatted_data = dict(data_source)  # Copy all existing data
        
        # Add some calculated/formatted fields
        formatted_data.update({
            'issue_date': datetime.now().strftime('%Y-%m-%d'),
            'expiry_date': StudentDataFormatter.calculate_expiry_date(
                data_source.get('admission_year', '') or data_source.get('admissionYear', '')
            ),
        })
        
        return formatted_data
    
    @staticmethod
    def calculate_expiry_date(admission_year):
        """
        Calculate ID card expiry date based on admission year.
        
        Args:
            admission_year (str/int): Year of admission
            
        Returns:
            str: Expiry date in YYYY-MM-DD format
        """
        if admission_year:
            try:
                # Assuming 4-year course, ID expires at graduation
                expiry_year = int(admission_year) + 4
                return f"{expiry_year}-12-31"
            except (ValueError, TypeError):
                pass
        
        # Default expiry: 2 years from now
        current_year = datetime.now().year
        return f"{current_year + 2}-12-31"